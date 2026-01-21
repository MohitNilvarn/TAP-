# File: app/services/document_processor.py
"""
Document processing service for extracting text from various file formats.
Supports PDF, DOCX, PPTX, TXT, and Markdown files.
"""
import os
from typing import List, Dict, Any, Optional
from abc import ABC, abstractmethod
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import markdown
from bs4 import BeautifulSoup

from app.core.logger import get_logger

logger = get_logger("TAP.DocumentProcessor")


class BaseDocumentProcessor(ABC):
    """Abstract base class for document processors."""
    
    @abstractmethod
    def extract_text(self, file_path: str) -> str:
        """Extract text content from the file."""
        pass
    
    @abstractmethod
    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from the file."""
        pass


class PDFProcessor(BaseDocumentProcessor):
    """Process PDF files using pypdf."""
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from all pages of a PDF."""
        try:
            reader = PdfReader(file_path)
            text_parts = []
            
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            raise
    
    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract PDF metadata."""
        try:
            reader = PdfReader(file_path)
            metadata = reader.metadata or {}
            
            return {
                "page_count": len(reader.pages),
                "title": getattr(metadata, "title", None),
                "author": getattr(metadata, "author", None),
                "subject": getattr(metadata, "subject", None),
                "creator": getattr(metadata, "creator", None),
            }
        except Exception as e:
            logger.error(f"Error extracting PDF metadata: {str(e)}")
            return {}


class DOCXProcessor(BaseDocumentProcessor):
    """Process DOCX files using python-docx."""
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from a DOCX file."""
        try:
            doc = DocxDocument(file_path)
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {str(e)}")
            raise
    
    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract DOCX metadata."""
        try:
            doc = DocxDocument(file_path)
            props = doc.core_properties
            
            return {
                "title": props.title,
                "author": props.author,
                "subject": props.subject,
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
            }
        except Exception as e:
            logger.error(f"Error extracting DOCX metadata: {str(e)}")
            return {}


class PPTXProcessor(BaseDocumentProcessor):
    """Process PPTX files using python-pptx."""
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from all slides of a PPTX."""
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = [f"--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                
                if len(slide_texts) > 1:  # More than just the slide header
                    text_parts.append("\n".join(slide_texts))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {str(e)}")
            raise
    
    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract PPTX metadata."""
        try:
            prs = Presentation(file_path)
            props = prs.core_properties
            
            return {
                "title": props.title,
                "author": props.author,
                "subject": props.subject,
                "slide_count": len(prs.slides),
            }
        except Exception as e:
            logger.error(f"Error extracting PPTX metadata: {str(e)}")
            return {}


class TextProcessor(BaseDocumentProcessor):
    """Process plain text and markdown files."""
    
    def extract_text(self, file_path: str) -> str:
        """Read text file content."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            
            # If markdown, convert to plain text
            if file_path.endswith(".md"):
                html = markdown.markdown(content)
                soup = BeautifulSoup(html, "html.parser")
                content = soup.get_text()
            
            return content
        except Exception as e:
            logger.error(f"Error reading text file: {str(e)}")
            raise
    
    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract text file metadata."""
        try:
            stat = os.stat(file_path)
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            
            return {
                "size_bytes": stat.st_size,
                "line_count": len(content.split("\n")),
                "word_count": len(content.split()),
            }
        except Exception as e:
            logger.error(f"Error extracting text metadata: {str(e)}")
            return {}


class DocumentProcessorFactory:
    """Factory for creating appropriate document processors."""
    
    _processors = {
        "pdf": PDFProcessor,
        "docx": DOCXProcessor,
        "pptx": PPTXProcessor,
        "txt": TextProcessor,
        "md": TextProcessor,
    }
    
    @classmethod
    def get_processor(cls, file_type: str) -> BaseDocumentProcessor:
        """Get the appropriate processor for a file type."""
        processor_class = cls._processors.get(file_type.lower())
        
        if not processor_class:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        return processor_class()
    
    @classmethod
    def supported_types(cls) -> List[str]:
        """Get list of supported file types."""
        return list(cls._processors.keys())


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """
    Split text into overlapping chunks for embedding.
    
    Args:
        text: The text to chunk
        chunk_size: Target size of each chunk in characters
        overlap: Number of characters to overlap between chunks
    
    Returns:
        List of text chunks
    """
    if not text:
        return []
    
    chunks = []
    start = 0
    text_len = len(text)
    
    while start < text_len:
        end = start + chunk_size
        
        # Try to break at a sentence or paragraph boundary
        if end < text_len:
            # Look for paragraph break
            para_break = text.rfind("\n\n", start, end)
            if para_break > start + chunk_size // 2:
                end = para_break
            else:
                # Look for sentence break
                for punct in [".", "!", "?"]:
                    sent_break = text.rfind(punct, start + chunk_size // 2, end)
                    if sent_break > 0:
                        end = sent_break + 1
                        break
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        start = end - overlap if end < text_len else text_len
    
    return chunks


def process_document(file_path: str, file_type: str) -> Dict[str, Any]:
    """
    Process a document and extract content.
    
    Args:
        file_path: Path to the file
        file_type: Type of file (pdf, docx, pptx, txt, md)
    
    Returns:
        Dictionary with extracted content and metadata
    """
    processor = DocumentProcessorFactory.get_processor(file_type)
    
    text = processor.extract_text(file_path)
    metadata = processor.extract_metadata(file_path)
    chunks = chunk_text(text)
    
    return {
        "content": text,
        "chunks": chunks,
        "metadata": metadata,
        "chunk_count": len(chunks),
    }
